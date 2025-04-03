import os
import logging
import tempfile
import streamlit as st
from dotenv import load_dotenv
from typing import List, Optional, Dict
from collections import Counter
import json
import spacy
from spacy.matcher import Matcher
import google.generativeai as genai
import pdfplumber
import pptx
import tempfile
import docx
from transformers import BartForConditionalGeneration, BartTokenizer
import torch
import re
from yt_dlp import YoutubeDL
from langchain_community.tools import YouTubeSearchTool
from langchain_community.tools import DuckDuckGoSearchResults
import streamlit.components.v1 as components
import ast
import warnings

# import subprocess
# subprocess.run(["python", "-m", "spacy", "download", "en_core_web_sm"])

# Ignore all warnings
warnings.filterwarnings('ignore')

# Set page configuration (must be the first Streamlit command)
st.set_page_config(
        page_title="Lecture Notes Assistant Pro",
        page_icon="📚",
        layout="wide",
        initial_sidebar_state="expanded"
    )

# Load environment variables
load_dotenv()

class DocumentProcessor:
    def __init__(self, api_key: str, gemini_model_name: str = "gemini-1.5-pro", gemini_flash_name: str = "gemini-1.5-flash"):
            """
            Initialize the document processor with optimized BART summarization and Gemini integration.
            """
            # Configure logging
            logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
            self.logger = logging.getLogger(__name__)

            # API Key Configuration for Gemini
            if not api_key:
                self.logger.error("API key is required for Gemini.")
                raise ValueError("API key not provided. Please input your Gemini API key.")

            try:
                # Configure Google AI (Gemini)
                genai.configure(api_key=api_key)
                self.gemini_model = genai.GenerativeModel(gemini_model_name)
                self.gemini_flash_model = genai.GenerativeModel(gemini_flash_name)
                
                # Initialize spaCy for NER
                try:
                    self.nlp = spacy.load("en_core_web_sm")
                    self.logger.info("spaCy model loaded successfully.")
                except Exception as e:
                    self.logger.warning(f"spaCy model loading failed: {e}")
                    self.nlp = None

                # Initialize dynamic NER matcher
                self.matcher = Matcher(self.nlp.vocab) if self.nlp else None
                self.conversation_history = []

                # Initialize BART for HF summarization
                try:
                    self.bart_model = BartForConditionalGeneration.from_pretrained('facebook/bart-large-cnn')
                    self.bart_tokenizer = BartTokenizer.from_pretrained('facebook/bart-large-cnn')
                    
                    # Enable CUDA if available for faster processing
                    if torch.cuda.is_available():
                        self.bart_model = self.bart_model.to('cuda')
                        self.logger.info("BART model loaded on GPU for faster processing.")
                    else:
                        self.logger.info("BART model loaded on CPU.")
                except Exception as e:
                    self.logger.error(f"BART model loading failed: {e}")
                    raise

            except Exception as e:
                self.logger.error(f"Initialization error: {e}")
                raise

    def extract_text(self, file_path: str) -> Optional[str]:
        """
        Extract text from various file formats as a single string.
        
        Returns:
        - A single string containing all the document text
        """
        if not os.path.exists(file_path):
            self.logger.error(f"File not found: {file_path}")
            return None

        ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if ext == ".pdf":
                with pdfplumber.open(file_path) as pdf:
                    return "\n\n".join([page.extract_text() or "" for page in pdf.pages])
            
            elif ext == ".pptx":
                prs = pptx.Presentation(file_path)
                return "\n\n".join([
                    "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
                    for slide in prs.slides
                ])
            
            elif ext in [".doc", ".docx"]:
                doc = docx.Document(file_path)
                return "\n".join([para.text for para in doc.paragraphs])
            
            elif ext == ".txt":
                with open(file_path, "r", encoding="utf-8") as f:
                    return f.read()
            
            else:
                self.logger.warning(f"Unsupported file type: {ext}")
                return None
        
        except Exception as e:
            self.logger.error(f"Error processing {file_path}: {e}")
            return None

    def setup_dynamic_ner(self, topic: str, document_text: str):
        """Efficiently sets up NER patterns using only document text and topic words."""
        if not self.matcher or not topic or not document_text:
            return
        
        self.matcher.remove("dynamic_entities", silent=True)

        # Step 1: Extract important words from the document using spaCy
        doc = self.nlp(document_text[:50000])  # Limit to first 50K chars for performance
        doc_keywords = [token.text.lower() for token in doc if token.pos_ in {"NOUN", "PROPN"}]

        # Step 2: Find most frequent words in the document
        word_freq = Counter(doc_keywords)
        top_keywords = [word for word, freq in word_freq.most_common(20)]  # Keep top 20 words

        # Step 3: Add topic words separately (splitting multi-word topics)
        topic_words = topic.lower().split()
        
        # Step 4: Combine both sets
        all_keywords = set(top_keywords + topic_words)

        # Step 5: Convert words into Matcher patterns
        patterns = [[{"LOWER": word}] for word in all_keywords]
        self.matcher.add("dynamic_entities", patterns)

        self.logger.info(f"Dynamic NER set up with {len(all_keywords)} keywords.")

    def split_text_into_pieces(self, text: str, max_tokens: int = 1024, overlap_percent: int = 5) -> List[str]:
        """Split text into chunks with minimal overlap for faster processing."""
        # Tokenize more efficiently
        tokens = self.bart_tokenizer(text, return_tensors="pt", truncation=False, add_special_tokens=False)["input_ids"][0]
        
        overlap_tokens = int(max_tokens * overlap_percent / 100)
        
        # Create chunks with minimal overlap
        chunks = []
        for i in range(0, len(tokens), max_tokens - overlap_tokens):
            chunk = tokens[i:i + max_tokens]
            chunks.append(chunk)
            
        # Convert token chunks back to text
        return [self.bart_tokenizer.decode(chunk, skip_special_tokens=True) for chunk in chunks]

    def bart_summarize_chunk(self, text: str, max_summary_length: int = 150) -> str:
        """Summarize a single chunk of text using BART with optimized parameters."""
        try:
            # Encode with proper truncation
            inputs = self.bart_tokenizer(text, return_tensors="pt", max_length=1024, truncation=True)
            
            # Move to GPU if available
            if torch.cuda.is_available():
                inputs = {k: v.to('cuda') for k, v in inputs.items()}
            
            # Generate with optimized parameters for speed
            summary_ids = self.bart_model.generate(
                inputs["input_ids"],
                attention_mask=inputs["attention_mask"],
                max_length=max_summary_length,
                min_length=max(30, int(max_summary_length / 4)),
                length_penalty=2.0,
                num_beams=2,
                early_stopping=True
            )
            
            # Convert back to CPU if needed
            if torch.cuda.is_available():
                summary_ids = summary_ids.cpu()
                
            return self.bart_tokenizer.decode(summary_ids[0], skip_special_tokens=True)
        except Exception as e:
            self.logger.error(f"BART chunk summarization error: {e}")
            return "Error summarizing chunk."

    def bart_summarize(self, text: str, max_tokens_per_chunk: int = 1024, max_summary_length: int = 250) -> str:
        """Summarize the text by chunking and summarizing each piece with optimized BART."""
        if not text:
            return "No text available for summarization."

        try:
            # Split text into larger chunks with minimal overlap
            chunks = self.split_text_into_pieces(text, max_tokens=max_tokens_per_chunk, overlap_percent=5)
            self.logger.info(f"Text split into {len(chunks)} chunks for BART summarization.")

            # Process chunks in parallel for faster summarization
            from concurrent.futures import ThreadPoolExecutor
            with ThreadPoolExecutor() as executor:
                summaries = list(executor.map(
                    lambda chunk: self.bart_summarize_chunk(chunk, max_summary_length), 
                    chunks
                ))

            # Combine summaries
            return " ".join(summaries)
        except Exception as e:
            self.logger.error(f"BART summarization error: {e}")
            return "Could not generate BART summary."

    def adjust_summary_with_feedback(self, gemini_summary: str, feedback: str) -> str:
        """Adjust the Gemini summary based on user feedback."""
        try:
            prompt = f"""
            Original Summary of lecture notes:
            {gemini_summary}

            Student Feedback:
            {feedback}

            Please revise the summary to address this feedback. Make the summary more conversational, 
            as if a teaching assistant is explaining the key points to a student. Use clear explanations 
            with relevant examples where appropriate. Focus on making complex concepts easier to understand.
            """
            response = self.gemini_model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            self.logger.error(f"Error adjusting summary with feedback: {e}")
            return "Could not adjust the summary based on feedback."

    def gemini_summarize(self, full_text: str, bart_summary: str, topic: str) -> str:
        """
        Generate a summary using Gemini, incorporating full text and BART summary.
        """
        if not full_text:
            return "No text available for summarization."

        try:
            # Extract key sections for context (beginning, middle, and end)
            text_length = len(full_text)
            beginning = full_text[:min(3000, text_length // 3)]
            middle = full_text[text_length // 2 - 1500:text_length // 2 + 1500] if text_length > 3000 else ""
            end = full_text[-min(3000, text_length // 3):] if text_length > 6000 else ""
            
            self.setup_dynamic_ner(topic, full_text)
            
            prompt = f"""
            I'm analyzing lecture notes on the topic of {topic}. Help me create a comprehensive, student-friendly summary.

            Here are portions of the lecture notes:
            Beginning:
            {beginning}
            
            {"Middle:" if middle else ""}
            {middle}
            
            {"End:" if end else ""}
            {end}

            I also have an initial summary created by BART:
            {bart_summary}

            Create a conversational, easy-to-understand summary that would help a student grasp the key concepts.
            Focus on the main ideas, important definitions, and connections between concepts.
            Use a friendly, helpful tone as if you're a teaching assistant explaining the material.
            Include relevant examples to illustrate complex ideas.
            Structure the summary in a logical way that follows the flow of the lecture.
            """
            
            response = self.gemini_model.generate_content(prompt)
            return response.text.strip()
        
        except Exception as e:
            self.logger.error(f"Gemini summary generation error: {e}")
            return "Could not generate Gemini summary."
            
    def generate_categorized_faqs(self, full_text: str, gemini_summary: str, topic: str, num_faqs: int = 9) -> dict:
        """Generate FAQs categorized by difficulty level (Easy, Medium, Hard)"""
        if not full_text:
            return {"easy": [], "medium": [], "hard": []}
            
        try:
            # Extract key sections for context (beginning, middle, and end)
            text_length = len(full_text)
            beginning = full_text[:min(3000, text_length // 3)]
            middle = full_text[text_length // 2 - 1500:text_length // 2 + 1500] if text_length > 3000 else ""
            end = full_text[-min(3000, text_length // 3):] if text_length > 6000 else ""
            
            prompt = f"""
            You are an expert in creating educational content from lecture notes. Based on the following lecture notes content and summary, generate {num_faqs} concise FAQs about '{topic}', evenly distributed across three difficulty levels (Easy, Medium, Hard). 

            Requirements:
            - Create exactly {num_faqs//3} questions for each difficulty level (Easy, Medium, Hard)
            - Easy questions should cover basic concepts and definitions
            - Medium questions should involve application of concepts or moderate understanding
            - Hard questions should require deep understanding, synthesis of multiple concepts, or critical thinking
            - Ensure each FAQ is unique and addresses common student questions
            - Make answers clear, concise, and suitable for students at the appropriate level
            - Frame questions as if they come from students attending a lecture

            Lecture Notes Content:
            Beginning:
            {beginning}
            
            {"Middle:" if middle else ""}
            {middle}
            
            {"End:" if end else ""}
            {end}
            
            Lecture Summary:
            {gemini_summary}

            Output the FAQs as a valid JSON object with the following structure:
            {{
              "easy": [
                {{ "question": "Student question here?", "answer": "Clear answer here." }},
                ...
              ],
              "medium": [
                {{ "question": "More complex student question here?", "answer": "More detailed answer here." }},
                ...
              ],
              "hard": [
                {{ "question": "Advanced student question here?", "answer": "Comprehensive answer here." }},
                ...
              ]
            }}

            Ensure the output is valid JSON that can be parsed with json.loads().
            """

            response = self.gemini_flash_model.generate_content(prompt)
            try:
                # Extract just the JSON part from the response
                json_text = response.text
                if "```json" in json_text:
                    json_text = json_text.split("```json")[1].split("```")[0].strip()
                elif "```" in json_text:
                    json_text = json_text.split("```")[1].strip()
                
                # Parse the JSON
                faqs_dict = json.loads(json_text)
                return faqs_dict
            except Exception as e:
                self.logger.error(f"Error parsing JSON from Gemini: {e}")
                return {
                    "easy": [{"question": "Error generating FAQs", "answer": "Please try again."}],
                    "medium": [],
                    "hard": []
                }
        except Exception as e:
            self.logger.error(f"FAQ generation error: {e}")
            return {
                "easy": [{"question": "Error generating FAQs", "answer": f"Error: {str(e)}"}],
                "medium": [],
                "hard": []
            }

    def answer_question(self, full_text: str, bart_summary: str, gemini_summary: str, topic: str, faqs_dict: dict, question: str) -> str:
        """
        Answer questions using full text context, BART summaries, Gemini summary, and FAQs.
        """
        if not full_text:
            return "No document text available for answering questions."

        try:
            # Add context from last 3 interactions
            history_context = "\n".join(f"Q: {q}\nA: {a}" for q, a in self.conversation_history[-3:])

            # Extract relevant sections based on question
            text_length = len(full_text)
            beginning = full_text[:min(2000, text_length // 4)]
            end = full_text[-min(2000, text_length // 4):] if text_length > 4000 else ""
            
            # Get middle sections that might be relevant
            middle_portion = ""
            if text_length > 6000 and self.nlp:
                # Look for keywords from the question in the text
                question_keywords = [token.text.lower() for token in self.nlp(question) if token.pos_ in {"NOUN", "VERB", "PROPN", "ADJ"}]
                
                # Split the document into sections and check for keyword matches
                sections = [full_text[i:i + 2000] for i in range(2000, text_length - 2000, 2000)]
                for section in sections[:3]:  # Limit to 3 sections for performance
                    for keyword in question_keywords:
                        if keyword in section.lower():
                            middle_portion += section + "\n\n"
                            break
            
            # Check if the question is similar to any FAQ
            faq_matches = []
            for difficulty in ["easy", "medium", "hard"]:
                if difficulty in faqs_dict:
                    for faq in faqs_dict[difficulty]:
                        # Simple keyword matching for now
                        q_words = set(question.lower().split())
                        faq_words = set(faq["question"].lower().split())
                        common_words = q_words.intersection(faq_words)
                        # If there's significant overlap, include this FAQ
                        if len(common_words) >= min(3, len(q_words) // 2):
                            faq_matches.append(f"Q: {faq['question']}\nA: {faq['answer']}")
            
            faq_context = "\n\n".join(faq_matches[:2])  # Use top 2 matching FAQs

            prompt = f"""
            You are a helpful teaching assistant for a course on {topic}. A student is asking you a question about the lecture notes.
            
            Here are relevant sections from the lecture notes:
            Beginning:
            {beginning}
            
            {"Relevant middle sections:" if middle_portion else ""}
            {middle_portion}
            
            {"End:" if end else ""}
            {end}

            Here's a summary of the lecture:
            {gemini_summary}
            
            {"Here are some similar questions that have been asked before:" if faq_context else ""}
            {faq_context}

            {"Previous questions from this student:" if history_context else ""}
            {history_context}

            The student asks: {question}

            Answer in a friendly, conversational, and helpful way, as if you're chatting with the student after class.
            Use examples or analogies where appropriate to illustrate concepts.
            If you don't have enough information from the lecture notes to answer confidently, be honest about it,
            but try to provide related information that might be helpful.
            Focus on explaining concepts clearly rather than just stating facts.
            """
            
            response = self.gemini_model.generate_content(prompt)
            answer = response.text.strip()
            
            # Store conversation history
            self.conversation_history.append((question, answer))
            
            return answer
        
        except Exception as e:
            self.logger.error(f"Question answering error: {e}")
            return "Could not answer the question."

    def extract_json(self, text: str) -> str:
        """Extracts JSON from Gemini's response if it contains extra text."""
        match = re.search(r"\[\s*\{.*\}\s*\]", text, re.DOTALL)
        return match.group(0) if match else None

    def generate_exam(self, topic: str, summary: str, difficulty: str = "medium") -> List[Dict[str, str]]:
        """Generate MCQs based on the document summary, topic, and difficulty."""
        prompt = f"""
        Create 10-15 multiple-choice questions based on the following topic and summary.
        Each question should have 4 options, with one correct answer clearly marked.
        The difficulty level of the questions should be {difficulty}.

        Format the output strictly as a valid JSON array:
        [
            {{
                "question": "What is the capital of France?",
                "options": ["Berlin", "Madrid", "Paris", "Rome"],
                "correct_answer": "Paris"
            }},
            ...
        ]
        
        Topic: {topic}
        Summary:
        {summary}
        """
        response = self.gemini_model.generate_content(prompt)

        # Log raw response for debugging
        raw_text = response.text.strip()
        
        # Extract only valid JSON if extra text is present
        json_text = self.extract_json(raw_text)
        if not json_text:
            self.logger.error("No valid JSON found in response.")
            return []

        # Parse JSON safely
        try:
            return json.loads(json_text)
        except json.JSONDecodeError as e:
            self.logger.error(f"Failed to parse exam questions: {e}")
            return []
        

    def generate_flashcards(self, topic: str, text: str, number_of_cards: int = 10) -> List[Dict[str, str]]:
        """Generate flashcards with topics and explanations based on document content."""
        prompt = f"""
        Create {number_of_cards} educational flashcards based on the following document about {topic}.
        Each flashcard should have:
        1. A concise topic or concept as the "front" (1-5 words)
        2. A brief explanation as the "back" (2-4 lines maximum)
        
        Format the output strictly as a valid JSON array:
        [
            {{
                "front": "Photosynthesis",
                "back": "The process by which green plants and some other organisms use sunlight to synthesize nutrients from carbon dioxide and water. Generates oxygen as a byproduct and is essential for life on Earth."
            }},
            ...
        ]
        
        Document Content:
        {text[:5000]}  # Using first 5000 chars for context
        """
        
        self.logger.info(f"Generating {number_of_cards} flashcards on topic: {topic}")
        response = self.gemini_model.generate_content(prompt)
        
        # Log raw response for debugging
        raw_text = response.text.strip()
        
        # Extract only valid JSON if extra text is present
        json_text = self.extract_json(raw_text)
        if not json_text:
            self.logger.error("No valid JSON found in flashcard response.")
            return []
        
        # Parse JSON safely
        try:
            return json.loads(json_text)
        except json.JSONDecodeError as e:
            self.logger.error(f"Failed to parse flashcards: {e}")
            return []



def load_css():
    """Load custom CSS for styling the FAQ section"""
    st.markdown("""
    <style>
    .faq-container {
        margin-bottom: 10px;
        border-radius: 5px;
        padding: 5px;
    }
    .easy-container {
        background-color: rgba(144, 238, 144, 0.2);
        border-left: 4px solid #90ee90;
    }
    .medium-container {
        background-color: rgba(255, 222, 173, 0.2);
        border-left: 4px solid #ffdead;
    }
    .hard-container {
        background-color: rgba(255, 182, 193, 0.2);
        border-left: 4px solid #ffb6c1;
    }
    .difficulty-header {
        font-size: 24px;
        font-weight: bold;
        margin: 20px 0 10px 0;
        padding: 5px;
        border-radius: 5px;
    }
    .easy-header {
        background-color: rgba(144, 238, 144, 0.4);
    }
    .medium-header {
        background-color: rgba(255, 222, 173, 0.4);
    }
    .hard-header {
        background-color: rgba(255, 182, 193, 0.4);
    }
    </style>
    """, unsafe_allow_html=True)


def display_faqs(faqs_dict):
    """Display FAQs with dropdowns and difficulty categories"""
    load_css()
    
    difficulty_levels = [
        ("easy", "Easy 🟢", "easy-header"),
        ("medium", "Medium 🟠", "medium-header"),
        ("hard", "Hard 🔴", "hard-header")
    ]
    
    for level, title, header_class in difficulty_levels:
        st.markdown(f"<div class='difficulty-header {header_class}'>{title}</div>", unsafe_allow_html=True)
        
        if level in faqs_dict and faqs_dict[level]:
            for i, faq in enumerate(faqs_dict[level]):
                with st.expander(f"Q{i+1}: {faq['question']}"):
                    st.markdown(f"<div class='faq-container {level}-container'>{faq['answer']}</div>", unsafe_allow_html=True)
        else:
            st.info(f"No {level} questions generated.")



def save_uploaded_file(uploaded_file):
    """Save the uploaded file to a temporary location and return the path."""
    try:
        # Create a temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f".{uploaded_file.name.split('.')[-1]}")
        temp_file.write(uploaded_file.getvalue())
        temp_file.close()
        return temp_file.name
    except Exception as e:
        st.error(f"Error saving uploaded file: {e}")
        return None
    

def render_flashcards(flashcards: List[Dict[str, str]]):
    """Render interactive flashcards in Streamlit."""
    
    st.markdown("""
    <style>
    .flashcard {
        position: relative;
        width: 100%;
        height: 200px;
        perspective: 1000px;
        margin-bottom: 20px;
    }
    
    .flashcard-inner {
        position: relative;
        width: 100%;
        height: 100%;
        text-align: center;
        transition: transform 0.6s;
        transform-style: preserve-3d;
    }
    
    .flashcard-flipped .flashcard-inner {
        transform: rotateY(180deg);
    }
    
    .flashcard-front, .flashcard-back {
        position: absolute;
        width: 100%;
        height: 100%;
        -webkit-backface-visibility: hidden;
        backface-visibility: hidden;
        display: flex;
        align-items: center;
        justify-content: center;
        border-radius: 10px;
        padding: 20px;
        box-sizing: border-box;
    }
    
    .flashcard-front {
        background-color: #f0f2f6;
        color: #262730;
        font-size: 24px;
        font-weight: bold;
    }
    
    .flashcard-back {
        background-color: #262730;
        color: white;
        transform: rotateY(180deg);
        font-size: 16px;
        overflow-y: auto;
        text-align: left;
    }
    </style>
    """, unsafe_allow_html=True)
    
    if 'flipped_cards' not in st.session_state:
        st.session_state.flipped_cards = [False] * len(flashcards)
    
    # Controls
    st.subheader("Flashcards")
    col1, col2 = st.columns([1, 1])
    with col1:
        if st.button("Flip All"):
            st.session_state.flipped_cards = [not state for state in st.session_state.flipped_cards]
    with col2:
        if st.button("Reset All"):
            st.session_state.flipped_cards = [False] * len(flashcards)
    
    # Display flashcards in a 2-column grid
    for i in range(0, len(flashcards), 2):
        col1, col2 = st.columns([1, 1])
        
        # Left column
        with col1:
            if i < len(flashcards):
                card_id = f"card_{i}"
                flipped_class = "flashcard-flipped" if st.session_state.flipped_cards[i] else ""
                
                html_content = f"""
                <div class="flashcard {flipped_class}" id="{card_id}">
                    <div class="flashcard-inner">
                        <div class="flashcard-front">
                            {flashcards[i]['front']}
                        </div>
                        <div class="flashcard-back">
                            {flashcards[i]['back']}
                        </div>
                    </div>
                </div>
                """
                st.markdown(html_content, unsafe_allow_html=True)
                
                if st.button(f"Flip Card {i+1}", key=f"flip_{i}"):
                    st.session_state.flipped_cards[i] = not st.session_state.flipped_cards[i]
                    st.rerun()
        
        # Right column
        with col2:
            j = i + 1
            if j < len(flashcards):
                card_id = f"card_{j}"
                flipped_class = "flashcard-flipped" if st.session_state.flipped_cards[j] else ""
                
                html_content = f"""
                <div class="flashcard {flipped_class}" id="{card_id}">
                    <div class="flashcard-inner">
                        <div class="flashcard-front">
                            {flashcards[j]['front']}
                        </div>
                        <div class="flashcard-back">
                            {flashcards[j]['back']}
                        </div>
                    </div>
                </div>
                """
                st.markdown(html_content, unsafe_allow_html=True)
                
                if st.button(f"Flip Card {j+1}", key=f"flip_{j}"):
                    st.session_state.flipped_cards[j] = not st.session_state.flipped_cards[j]
                    st.rerun()

def render_flashcards_advanced(flashcards: List[Dict[str, str]]):
    """Render interactive flashcards with advanced CSS animations."""
    
    st.markdown("""
    <style>
    /* Flashcard container */
    .flashcard-container {
        display: grid;
        grid-template-columns: repeat(auto-fill, minmax(300px, 1fr));
        gap: 20px;
        margin-top: 20px;
    }
    
    /* Individual flashcard */
    .flashcard-adv {
        background-color: transparent;
        perspective: 1000px;
        height: 200px;
        cursor: pointer;
        margin-bottom: 20px;
    }
    
    /* Flashcard inner container for flip effect */
    .flashcard-inner-adv {
        position: relative;
        width: 100%;
        height: 100%;
        text-align: center;
        transition: transform 0.8s;
        transform-style: preserve-3d;
        box-shadow: 0 4px 8px 0 rgba(0,0,0,0.2);
        border-radius: 10px;
    }
    
    /* Flipped state */
    .flipped .flashcard-inner-adv {
        transform: rotateY(180deg);
    }
    
    /* Front and back sides */
    .flashcard-front-adv, .flashcard-back-adv {
        position: absolute;
        width: 100%;
        height: 100%;
        -webkit-backface-visibility: hidden;
        backface-visibility: hidden;
        display: flex;
        flex-direction: column;
        align-items: center;
        justify-content: center;
        padding: 20px;
        box-sizing: border-box;
        border-radius: 10px;
    }
    
    /* Front styling */
    .flashcard-front-adv {
        background: linear-gradient(135deg, #6e8efb, #a777e3);
        color: white;
    }
    
    /* Front title */
    .flashcard-title {
        font-size: 24px;
        font-weight: bold;
        text-shadow: 1px 1px 3px rgba(0,0,0,0.3);
    }
    
    /* Small prompt text */
    .flip-prompt {
        position: absolute;
        bottom: 10px;
        font-size: 12px;
        opacity: 0.7;
    }
    
    /* Back styling */
    .flashcard-back-adv {
        background: linear-gradient(135deg, #2b5876, #4e4376);
        color: white;
        transform: rotateY(180deg);
        text-align: left;
        overflow-y: auto;
    }
    
    /* Back content */
    .flashcard-content {
        font-size: 16px;
        line-height: 1.5;
    }
    
    /* Controls */
    .flashcard-controls {
        display: flex;
        justify-content: space-between;
        margin-bottom: 20px;
    }
    
    /* Control buttons */
    .control-btn {
        background-color: #4CAF50;
        border: none;
        color: white;
        padding: 10px 20px;
        text-align: center;
        text-decoration: none;
        display: inline-block;
        font-size: 16px;
        margin: 4px 2px;
        transition-duration: 0.4s;
        cursor: pointer;
        border-radius: 5px;
    }
    
    .control-btn:hover {
        background-color: white;
        color: black;
        border: 1px solid #4CAF50;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # Initialize flipped state in session state if not already present
    if 'flipped_states' not in st.session_state:
        st.session_state.flipped_states = [False] * len(flashcards)
    
    # Control buttons
    st.markdown('<div class="flashcard-controls">', unsafe_allow_html=True)
    col1, col2, col3, col4 = st.columns([1, 1, 1, 1])
    
    with col1:
        if st.button("Flip All Cards", key="flip_all"):
            # Toggle all cards
            all_flipped = all(st.session_state.flipped_states)
            st.session_state.flipped_states = [not all_flipped] * len(flashcards)
            st.rerun()
    
    with col2:
        if st.button("Reset All Cards", key="reset_all"):
            # Set all cards to front side
            st.session_state.flipped_states = [False] * len(flashcards)
            st.rerun()
    
    with col3:
        if st.button("Study Mode", key="study_mode"):
            # Start with all cards front-side up
            st.session_state.flipped_states = [False] * len(flashcards)
            st.session_state.study_mode = True
            st.rerun()
    
    with col4:
        if st.button("Export Cards", key="export_cards"):
            # Export functionality is handled separately
            pass
    
    st.markdown('</div>', unsafe_allow_html=True)
    
    # Start flashcard grid container
    st.markdown('<div class="flashcard-container">', unsafe_allow_html=True)
    
    # Create each flashcard
    for i, card in enumerate(flashcards):
        # Generate HTML for a single flashcard
        flipped_class = "flipped" if st.session_state.flipped_states[i] else ""
        
        # Format the back content to handle potential HTML
        back_content = card['back'].replace('\n', '<br>')
        
        flashcard_html = f"""
        <div class="flashcard-adv {flipped_class}" id="card_{i}" onclick="this.classList.toggle('flipped')">
            <div class="flashcard-inner-adv">
                <div class="flashcard-front-adv">
                    <div class="flashcard-title">{card['front']}</div>
                    <div class="flip-prompt">Click to flip</div>
                </div>
                <div class="flashcard-back-adv">
                    <div class="flashcard-content">{back_content}</div>
                </div>
            </div>
        </div>
        """
        
        # Use columns for buttons under each card
        st.markdown(flashcard_html, unsafe_allow_html=True)
        
        # Add a button to flip each card programmatically
        if st.button(f"Flip Card {i+1}", key=f"flip_btn_{i}"):
            st.session_state.flipped_states[i] = not st.session_state.flipped_states[i]
            st.rerun()
    
    # Close the flashcard container
    st.markdown('</div>', unsafe_allow_html=True)
    
    # Add instructions
    st.markdown("""
    **Instructions:**
    - Click directly on a card to flip it
    - Use the buttons to control individual cards or all cards at once
    - In Study Mode, go through the cards one by one
    """)


def get_video_details(video_url):
    """Extract title and thumbnail of a YouTube video using yt_dlp"""
    try:
        ydl_opts = {
            'quiet': True,
            'noplaylist': True,
            'skip_download': True,
            'no_warnings': True
        }
        with YoutubeDL(ydl_opts) as ydl:
            info_dict = ydl.extract_info(video_url, download=False)
            title = info_dict.get('title', 'Title not found')
            thumbnail = info_dict.get('thumbnail', None)
            duration = info_dict.get('duration')
            
            # Format duration as MM:SS
            if duration:
                minutes, seconds = divmod(duration, 60)
                duration_str = f"{minutes}:{seconds:02d}"
            else:
                duration_str = "Unknown"
                
            return {
                "title": title,
                "thumbnail": thumbnail,
                "duration": duration_str,
                "views": info_dict.get('view_count', 'Unknown'),
                "channel": info_dict.get('uploader', 'Unknown')
            }
    except Exception as e:
        st.error(f"Error getting video details: {e}")
        return {
            "title": "Error retrieving details",
            "thumbnail": None,
            "duration": "Unknown",
            "views": "Unknown",
            "channel": "Unknown"
        }

def search_youtube(query, num_results=5):
    """Search YouTube for videos related to the query"""
    with st.spinner(f"Searching for '{query}' on YouTube..."):
        tool = YouTubeSearchTool()
        results = tool.run(f"{query}, {num_results}")
        
        # Convert string representation of list to actual list
        try:
            # Method 1: Using ast.literal_eval for safe evaluation
            if results.strip().startswith('[') and results.strip().endswith(']'):
                url_list = ast.literal_eval(results.strip())
            # Method 2: Using regex as fallback
            else:
                url_list = re.findall(r'https://www\.youtube\.com/watch\?v=[^\'"\s]+', results)
            
            # Convert URLs to video objects with details
            videos = []
            for i, video_url in enumerate(url_list[:num_results]):
                with st.spinner(f"Getting details for video {i+1}/{min(len(url_list), num_results)}..."):
                    details = get_video_details(video_url)
                    videos.append({"url": video_url, **details})
            return videos
        except Exception as e:
            st.error(f"Error parsing YouTube results: {e}")
            st.code(results)  # Display the raw results for debugging
            return []

def get_study_links(topic, max_results=4):
    """Get educational reference links for a topic using DuckDuckGo"""
    with st.spinner(f"Finding reference materials for '{topic}'..."):
        try:
            search = DuckDuckGoSearchResults(output_format="list")
            prompt = f"Best educational resources to study {topic}. Provide academic and research-based links."
            results = search.invoke(prompt)
            
            # Process and clean up results (limit to max_results)
            clean_results = []
            for item in results[:max_results]:
                if isinstance(item, dict) and 'link' in item and 'title' in item and 'snippet' in item:
                    clean_results.append({
                        'link': item['link'],
                        'title': item['title'],
                        'snippet': item['snippet']
                    })
            return clean_results
        except Exception as e:
            st.error(f"Error getting reference links: {e}")
            return []

def extract_video_id(url):
    """Extract the video ID from a YouTube URL"""
    if "youtube.com/watch?v=" in url:
        video_id = url.split("youtube.com/watch?v=")[1]
        # Remove any query parameters after the ID
        if "&" in video_id:
            video_id = video_id.split("&")[0]
        return video_id
    elif "youtu.be/" in url:
        return url.split("youtu.be/")[1].split("?")[0]
    return None


# Initialize Streamlit app
def main():
    st.title("📚 Lecture Notes Assistant Pro")
    st.write("Upload your lecture notes to get a complete study package with summaries, FAQs, and a personalized AI study buddy!")

    # Sidebar for API key input
    with st.sidebar:
        st.header("🔑 API Key Configuration")
        api_key = st.text_input("Enter your Gemini API Key", type="password", placeholder="Your Gemini API Key")
        
        if not api_key:
            st.warning("Please enter your Gemini API key to proceed.")
            return

    # Initialize session state variables if they don't exist
    if 'processor' not in st.session_state:
        try:
            st.session_state.processor = DocumentProcessor(api_key=api_key)
            st.session_state.full_text = None
            st.session_state.bart_summary = None
            st.session_state.gemini_summary = None
            st.session_state.faqs_dict = None
            st.session_state.topic = None
            st.session_state.file_info = None
            st.session_state.conversation_history = []
            st.session_state.num_faqs_per_level = 3
        except Exception as e:
            st.error(f"Error initializing the assistant: {str(e)}")
            return

    # Create sidebar for file upload and topic input
    with st.sidebar:
        st.header("📝 Upload Lecture Notes")
        
        uploaded_file = st.file_uploader("Upload your notes", type=["pdf", "docx", "txt", "pptx"] ,  label_visibility="visible")
        
        if uploaded_file is not None:
            # Create a temporary file to process
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{uploaded_file.name.split('.')[-1]}") as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                temp_path = tmp_file.name
            
            # Process the file if it's new or different from the previous one
            if st.session_state.file_info != uploaded_file.name:
                with st.spinner("Reading your lecture notes..."):
                    st.session_state.full_text = st.session_state.processor.extract_text(temp_path)
                    st.session_state.file_info = uploaded_file.name
                    # Reset summaries and topic when a new file is uploaded
                    st.session_state.bart_summary = None
                    st.session_state.gemini_summary = None
                    st.session_state.faqs_dict = None
                    st.session_state.topic = None
                    st.session_state.processor.conversation_history = []
                    st.session_state.conversation_history = []
                    st.session_state.exam_questions = []
                    st.session_state.flashcards = []
                    st.session_state.user_answers = {}
                    st.session_state.score_calculated = False

                if st.session_state.full_text:
                    st.success(f"Notes successfully uploaded!")
                    st.info(f"Notes Stats:\n- Length: {len(st.session_state.full_text)} characters\n- Words: {len(st.session_state.full_text.split())} words")
                else:
                    st.error("Couldn't read the notes. Try a different file format.")
            
            # Remove the temporary file
            try:
                os.unlink(temp_path)
            except:
                pass
        
        # Topic input
        if st.session_state.full_text is not None:
            st.session_state.topic = st.text_input("What subject/course are these notes for?", value=st.session_state.topic if st.session_state.topic else "")
            
            if st.session_state.topic and not st.session_state.bart_summary:
                st.session_state.num_faqs_per_level = st.slider("Number of FAQs per difficulty level", 1, 5, 3)
                
                if st.button("Process Lecture Notes"):
                    # Generate BART summary
                    with st.spinner("Analyzing your notes..."):
                        st.session_state.bart_summary = st.session_state.processor.bart_summarize(st.session_state.full_text)
                    
                    # Generate Gemini summary
                    with st.spinner("Creating your personalized study guide..."):
                        st.session_state.gemini_summary = st.session_state.processor.gemini_summarize(
                            st.session_state.full_text, 
                            st.session_state.bart_summary, 
                            st.session_state.topic
                        )
                    
                    # Generate FAQs
                    with st.spinner("Generating frequently asked questions..."):
                        st.session_state.faqs_dict = st.session_state.processor.generate_categorized_faqs(
                            st.session_state.full_text,
                            st.session_state.gemini_summary,
                            st.session_state.topic,
                            num_faqs=st.session_state.num_faqs_per_level * 3
                        )
                    
                    st.success("Your study materials are ready!")

        if st.session_state.gemini_summary:
                st.header("Learning Tools")
                tool_type = st.selectbox(
                    "Select tool type:",
                    options=["Exam Questions", "Flashcards"],
                    index=0
                )
                
                difficulty = st.selectbox(
                    "Select difficulty level:",
                    options=["easy", "medium", "hard"],
                    index=1
                )
                
                num_items = st.slider("Number of items to generate:", min_value=5, max_value=20, value=10)
                
                if st.button("Generate Learning Material"):
                    if tool_type == "Exam Questions":
                        with st.spinner("Generating exam questions..."):
                            st.session_state.exam_questions = st.session_state.processor.generate_exam(
                                st.session_state.topic,
                                st.session_state.gemini_summary,
                                difficulty
                            )
                            # Reset user answers and score when generating new questions
                            st.session_state.user_answers = {}
                            st.session_state.score_calculated = False
                            
                            if st.session_state.exam_questions:
                                st.success(f"Generated {len(st.session_state.exam_questions)} questions!")
                            else:
                                st.error("Failed to generate questions. Please try again.")
                    
                    elif tool_type == "Flashcards":
                        with st.spinner("Generating flashcards..."):
                            st.session_state.flashcards = st.session_state.processor.generate_flashcards(
                                st.session_state.topic,
                                st.session_state.gemini_summary,
                                num_items
                            )
                            
                            if st.session_state.flashcards:
                                st.success(f"Generated {len(st.session_state.flashcards)} flashcards!")
                            else:
                                st.error("Failed to generate flashcards. Please try again.")


    # Main content area
    if st.session_state.full_text is not None:
        # Create tabs for different functionalities
        tab1, tab2, tab3, tab4 , tab5 , tab6 , tab7 = st.tabs(["Study Guide", "FAQ Generator", "Study Buddy Chat", "Exam Questions", "Flashcards" , "References and Videos" ,"Original Notes"])
         
        # Study Guide tab
        with tab1:
            if st.session_state.bart_summary and st.session_state.gemini_summary:
                st.subheader("Your Personalized Study Guide")
                st.write(st.session_state.gemini_summary)
                
                # Summary feedback
                st.divider()
                st.subheader("How can we improve your study guide?")
                feedback = st.text_area("Let us know what you'd like to change:", height=100, 
                                       placeholder="Example: 'Add more examples about topic X' or 'Make it more concise'")
                
                if st.button("Improve My Study Guide"):
                    with st.spinner("Updating your study guide..."):
                        st.session_state.gemini_summary = st.session_state.processor.adjust_summary_with_feedback(
                            st.session_state.gemini_summary, 
                            feedback
                        )
                    st.success("Study guide updated based on your feedback!")
                    st.rerun()
            else:
                if st.session_state.topic:
                    st.info("Click 'Process Lecture Notes' in the sidebar to generate your personalized study materials.")
                else:
                    st.info("Please enter the subject or course name in the sidebar.")
        
        # FAQ Generator tab
        with tab2:
            if st.session_state.faqs_dict:
                st.subheader("Frequently Asked Questions")
                display_faqs(st.session_state.faqs_dict)
                
                # Add download button for FAQs
                st.divider()
                
                # Convert FAQs to markdown for download
                markdown_faqs = f"# {st.session_state.topic} FAQs\n\n"
                for level, title in [("easy", "Easy Questions"), ("medium", "Medium Questions"), ("hard", "Hard Questions")]:
                    markdown_faqs += f"## {title}\n\n"
                    if level in st.session_state.faqs_dict:
                        for i, faq in enumerate(st.session_state.faqs_dict[level]):
                            markdown_faqs += f"### Q{i+1}: {faq['question']}\n\n{faq['answer']}\n\n"
                
                col1, col2 = st.columns(2)
                with col1:
                    st.download_button(
                        label="Download FAQs as Markdown",
                        data=markdown_faqs,
                        file_name=f"{st.session_state.topic}_FAQs.md",
                        mime="text/markdown"
                    )
                
                with col2:
                    if st.button("Regenerate FAQs"):
                        with st.spinner("Regenerating FAQs..."):
                            st.session_state.faqs_dict = st.session_state.processor.generate_categorized_faqs(
                                st.session_state.full_text,
                                st.session_state.gemini_summary,
                                st.session_state.topic,
                                num_faqs=st.session_state.num_faqs_per_level * 3
                            )
                        st.success("FAQs regenerated!")
                        st.rerun()
                
                # Option to regenerate specific difficulty level
                st.subheader("Regenerate Specific Category")
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    if st.button("Regenerate Easy FAQs"):
                        with st.spinner("Regenerating easy questions..."):
                            new_prompt = f"""
                            Based on the lecture notes about '{st.session_state.topic}', generate {st.session_state.num_faqs_per_level} EASY questions and answers.
                            These should cover basic concepts and definitions. Output as valid JSON array of objects with 'question' and 'answer' fields.
                            """
                            response = st.session_state.processor.gemini_flash_model.generate_content(new_prompt)
                            try:
                                json_text = response.text
                                if "```json" in json_text:
                                    json_text = json_text.split("```json")[1].split("```")[0].strip()
                                elif "```" in json_text:
                                    json_text = json_text.split("```")[1].strip()
                                
                                new_easy = json.loads(json_text)
                                if not isinstance(new_easy, list):
                                    # Handle case where output might be a dictionary with a key like "questions"
                                    for key in new_easy:
                                        if isinstance(new_easy[key], list):
                                            new_easy = new_easy[key]
                                            break
                                
                                st.session_state.faqs_dict["easy"] = new_easy
                                st.rerun()
                            except Exception as e:
                                st.error(f"Error regenerating easy FAQs: {e}")

                with col2:
                                    if st.button("Regenerate Medium FAQs"):
                                        with st.spinner("Regenerating medium questions..."):
                                            new_prompt = f"""
                                            Based on the lecture notes about '{st.session_state.topic}', generate {st.session_state.num_faqs_per_level} MEDIUM difficulty questions and answers.
                                            These should involve application of concepts or moderate understanding. Output as valid JSON array of objects with 'question' and 'answer' fields.
                                            """
                                            response = st.session_state.processor.gemini_flash_model.generate_content(new_prompt)
                                            try:
                                                json_text = response.text
                                                if "```json" in json_text:
                                                    json_text = json_text.split("```json")[1].split("```")[0].strip()
                                                elif "```" in json_text:
                                                    json_text = json_text.split("```")[1].strip()
                                                
                                                new_medium = json.loads(json_text)
                                                if not isinstance(new_medium, list):
                                                    for key in new_medium:
                                                        if isinstance(new_medium[key], list):
                                                            new_medium = new_medium[key]
                                                            break
                                                
                                                st.session_state.faqs_dict["medium"] = new_medium
                                                st.rerun()
                                            except Exception as e:
                                                st.error(f"Error regenerating medium FAQs: {e}")
                                
                with col3:
                                    if st.button("Regenerate Hard FAQs"):
                                        with st.spinner("Regenerating hard questions..."):
                                            new_prompt = f"""
                                            Based on the lecture notes about '{st.session_state.topic}', generate {st.session_state.num_faqs_per_level} HARD questions and answers.
                                            These should require deep understanding, synthesis of multiple concepts, or critical thinking. Output as valid JSON array of objects with 'question' and 'answer' fields.
                                            """
                                            response = st.session_state.processor.gemini_flash_model.generate_content(new_prompt)
                                            try:
                                                json_text = response.text
                                                if "```json" in json_text:
                                                    json_text = json_text.split("```json")[1].split("```")[0].strip()
                                                elif "```" in json_text:
                                                    json_text = json_text.split("```")[1].strip()
                                                
                                                new_hard = json.loads(json_text)
                                                if not isinstance(new_hard, list):
                                                    for key in new_hard:
                                                        if isinstance(new_hard[key], list):
                                                            new_hard = new_hard[key]
                                                            break
                                                
                                                st.session_state.faqs_dict["hard"] = new_hard
                                                st.rerun()
                                            except Exception as e:
                                                st.error(f"Error regenerating hard FAQs: {e}")
            else:
                st.info("Generate your study guide first to see categorized FAQs.")
                        
        # Study Buddy Chat tab
        with tab3:
            if st.session_state.bart_summary and st.session_state.gemini_summary and st.session_state.faqs_dict:
                st.subheader("Ask Your AI Study Buddy")
                                
                # Create a chat interface
                for i, (q, a) in enumerate(st.session_state.conversation_history):
                    with st.chat_message("user"):
                        st.write(q)
                    with st.chat_message("assistant", avatar="📚"):
                        st.write(a)
                                
                # Question input using chat_input for better UI
                question = st.chat_input("Ask a question about your lecture notes...")
                                
                if question:
                    # Show user question
                    with st.chat_message("user"):
                        st.write(question)
                                    
                    # Generate and show answer
                    with st.chat_message("assistant", avatar="📚"):
                        with st.spinner("Thinking..."):
                            answer = st.session_state.processor.answer_question(
                                st.session_state.full_text,
                                st.session_state.bart_summary,
                                st.session_state.gemini_summary,
                                st.session_state.topic,
                                st.session_state.faqs_dict,
                                question
                                )
                        st.write(answer)
                                    
                    # Store in session state for display persistence
                    st.session_state.conversation_history.append((question, answer))
                    st.rerun()

            else:
                st.info("Generate your study guide and FAQs first before chatting with your AI study buddy.")



        #Exam Questions tab           
        with tab4:
            st.header("Exam Questions")
            
            if st.session_state.exam_questions:
                # Create a form for the quiz
                with st.form("quiz_form"):
                    for i, q in enumerate(st.session_state.exam_questions):
                        st.markdown(f"### Question {i+1}: {q['question']}")
                        
                        # Generate a unique key for each question
                        question_key = f"q_{i}"
                        
                        # Use radio buttons for answer selection
                        selected_option = st.radio(
                            f"Select your answer for question {i+1}:",
                            options=q['options'],
                            key=question_key,
                            index=None  # No default selection
                        )
                        
                        # Store the selected answer in session state
                        if selected_option is not None:
                            st.session_state.user_answers[question_key] = selected_option
                        
                        st.markdown("---")
                    
                    # Submit button for the form
                    submitted = st.form_submit_button("Check Score")
                    
                    if submitted:
                        st.session_state.score_calculated = True
                        
                        # Calculate score
                        total_questions = len(st.session_state.exam_questions)
                        correct_answers = 0
                        
                        for i, q in enumerate(st.session_state.exam_questions):
                            question_key = f"q_{i}"
                            if question_key in st.session_state.user_answers:
                                user_answer = st.session_state.user_answers[question_key]
                                if user_answer == q['correct_answer']:
                                    correct_answers += 1
                        
                        st.session_state.score = correct_answers
                
                # Display results if score has been calculated
                if st.session_state.score_calculated:
                    total_questions = len(st.session_state.exam_questions)
                    score_percentage = (st.session_state.score / total_questions * 100) if total_questions > 0 else 0
                    
                    st.success(f"Your Score: {st.session_state.score}/{total_questions} ({score_percentage:.1f}%)")
                    
                    # Display detailed results
                    st.subheader("Results")
                    for i, q in enumerate(st.session_state.exam_questions):
                        question_key = f"q_{i}"
                        user_answer = st.session_state.user_answers.get(question_key, "Not answered")
                        correct_answer = q['correct_answer']
                        
                        st.markdown(f"**Question {i+1}:** {q['question']}")
                        
                        if user_answer == correct_answer:
                            st.markdown(f"✅ Your answer: **{user_answer}** (Correct)")
                        else:
                            st.markdown(f"❌ Your answer: **{user_answer}**")
                            st.markdown(f"✅ Correct answer: **{correct_answer}**")
                        
                        st.markdown("---")
                
                # Export options
                col1, col2 = st.columns(2)
                with col1:
                    if st.button("Export Questions as JSON"):
                        json_str = json.dumps(st.session_state.exam_questions, indent=2)
                        st.download_button(
                            label="Download JSON",
                            data=json_str,
                            file_name=f"{st.session_state.topic}_exam_questions.json",
                            mime="application/json",
                            key="download_json"
                        )
            else:
                st.info("No exam questions generated yet. Use the sidebar to generate questions.")

        #Flashcards
        with tab5:
            st.header("Flashcards")
            
            if st.session_state.flashcards:
                # Render flashcards using the dedicated function
                render_flashcards(st.session_state.flashcards)
                
                # Export options
                if st.button("Export Flashcards as JSON"):
                    json_str = json.dumps(st.session_state.flashcards, indent=2)
                    st.download_button(
                        label="Download JSON",
                        data=json_str,
                        file_name=f"{st.session_state.topic}_flashcards.json",
                        mime="application/json",
                        key="download_flashcards_json"
                    )
            else:
                st.info("No flashcards generated yet. Use the sidebar to generate flashcards.")


        #Videos and References
        with tab6:
            # Custom CSS for better styling
            st.markdown("""
            <style>
                .video-card {
                    border-radius: 10px;
                    border: 1px solid #ddd;
                    padding: 10px;
                    margin-bottom: 20px;
                    transition: transform 0.3s;
                    height: 100%;
                }
                .video-card:hover {
                    transform: translateY(-5px);
                    box-shadow: 0 10px 20px rgba(0,0,0,0.1);
                }
                .reference-card {
                    border-radius: 10px;
                    border: 1px solid #e0e0ff;
                    background-color: #f8f9ff;
                    padding: 15px;
                    margin-bottom: 15px;
                    transition: transform 0.2s;
                }
                .reference-card:hover {
                    transform: translateY(-3px);
                    box-shadow: 0 5px 15px rgba(0,0,0,0.08);
                }
                .search-container {
                    background-color: #f9f9f9;
                    padding: 20px;
                    border-radius: 10px;
                    margin-bottom: 20px;
                }
                .player-container {
                    background-color: #f0f0f0;
                    padding: 20px;
                    border-radius: 10px;
                    margin-top: 30px;
                }
                .title {
                    color: #1E3A8A;
                    font-size: 36px;
                    font-weight: bold;
                    margin-bottom: 30px;
                    text-align: center;
                }
                .subtitle {
                    color: #1E3A8A;
                    border-left: 4px solid #1E3A8A;
                    padding-left: 10px;
                }
                .meta-info {
                    color: #666;
                    font-size: 14px;
                }
                .tab-content {
                    padding: 20px 0;
                }
                .reference-title {
                    color: #2C5282;
                    font-weight: bold;
                    margin-bottom: 5px;
                }
                .reference-snippet {
                    color: #4A5568;
                    font-size: 14px;
                    margin-bottom: 10px;
                }
                .reference-link {
                    color: #3182CE;
                    font-size: 14px;
                    word-break: break-word;
                }
            </style>
            """, unsafe_allow_html=True)
            # Page title with icon
            st.markdown('<p class="title">🎬 Learning Resource Explorer</p>', unsafe_allow_html=True)

            # Search form in a container
            st.markdown('<div class="search-container">', unsafe_allow_html=True)
            col1, col2 = st.columns([3, 1])
            with col1:
                search_query = st.text_input("What would you like to learn about today?", 
                                            st.session_state.topic,
                                            placeholder="Enter a topic to search")
            with col2:
                num_results = st.slider("Number of videos", min_value=1, max_value=10, value=5)

            if st.button("🔍 Search for Learning Resources", use_container_width=True):
                # Store the search query
                st.session_state.search_query = search_query
                
                # Search for videos
                st.session_state.videos = search_youtube(search_query, num_results)
                
                # Search for reference links
                st.session_state.references = get_study_links(search_query)
                
                # Reset selected video
                st.session_state.selected_video = None
            st.markdown('</div>', unsafe_allow_html=True)

            # Initialize session states if they don't exist
            if 'videos' not in st.session_state:
                st.session_state.videos = []
            if 'references' not in st.session_state:
                st.session_state.references = []
            if 'selected_video' not in st.session_state:
                st.session_state.selected_video = None
            if 'search_query' not in st.session_state:
                st.session_state.search_query = ""

            # Create tabs for Videos and References
            if st.session_state.videos or st.session_state.references:
                tab1, tab2 = st.tabs(["📺 Video Tutorials", "📚 Reference Materials"])
                
                # Videos Tab
                with tab1:
                    st.markdown('<div class="tab-content">', unsafe_allow_html=True)
                    if st.session_state.videos:
                        # Display videos in a grid
                        cols = st.columns(min(3, len(st.session_state.videos)))
                        for i, video in enumerate(st.session_state.videos):
                            col_index = i % len(cols)
                            with cols[col_index]:
                                st.markdown(f'<div class="video-card">', unsafe_allow_html=True)
                                
                                # Display thumbnail with play button overlay
                                if video["thumbnail"]:
                                    st.image(video["thumbnail"], use_container_width=True)
                                else:
                                    # Fallback if no thumbnail
                                    st.image("https://via.placeholder.com/320x180?text=No+Thumbnail", use_column_width=True)
                                
                                # Display video title and metadata
                                st.markdown(f"**{video['title']}**")
                                st.markdown(f'<p class="meta-info">👤 {video["channel"]} | ⏱️ {video["duration"]} | 👁️ {video["views"]:,}</p>', 
                                        unsafe_allow_html=True)
                                
                                # Play button
                                if st.button("▶️ Play", key=f"play_{i}", use_container_width=True):
                                    st.session_state.selected_video = video
                                    st.rerun()
                                    
                                st.markdown('</div>', unsafe_allow_html=True)
                    else:
                        st.info("No videos found. Try a different search query.")
                    st.markdown('</div>', unsafe_allow_html=True)
                
                # References Tab
                with tab2:
                    st.markdown('<div class="tab-content">', unsafe_allow_html=True)
                    if st.session_state.references:
                        for i, ref in enumerate(st.session_state.references):
                            st.markdown(f'<div class="reference-card">', unsafe_allow_html=True)
                            st.markdown(f'<div class="reference-title">{ref["title"]}</div>', unsafe_allow_html=True)
                            st.markdown(f'<div class="reference-snippet">{ref["snippet"]}</div>', unsafe_allow_html=True)
                            st.markdown(f'<a href="{ref["link"]}" target="_blank" class="reference-link">{ref["link"]}</a>', unsafe_allow_html=True)
                            st.markdown('</div>', unsafe_allow_html=True)
                    else:
                        st.info("No reference materials found. Try a different search query.")
                    st.markdown('</div>', unsafe_allow_html=True)

            # Display selected video player
            if st.session_state.selected_video:
                st.markdown('<div class="player-container">', unsafe_allow_html=True)
                st.markdown('<h2 class="subtitle">Now Playing</h2>', unsafe_allow_html=True)
                
                video = st.session_state.selected_video
                
                col1, col2 = st.columns([2, 1])
                
                with col1:
                    video_id = extract_video_id(video['url'])
                    if video_id:
                        # Create an iframe to embed the YouTube video
                        embed_code = f'''
                        <div style="position: relative; padding-bottom: 56.25%; height: 0; overflow: hidden; max-width: 100%; border-radius: 10px;">
                            <iframe 
                                style="position: absolute; top: 0; left: 0; width: 100%; height: 90%; border-radius: 10px;" 
                                src="https://www.youtube.com/embed/{video_id}" 
                                frameborder="0" 
                                allow="accelerometer; autoplay; clipboard-write; encrypted-media; gyroscope; picture-in-picture" 
                                allowfullscreen>
                            </iframe>
                        </div>
                        '''
                        components.html(embed_code, height=400)
                    else:
                        st.error("Could not extract video ID from the URL")
                    
                    # Show references relevant to the current video
                    if st.session_state.references:
                        st.markdown("### 📚 Reference Materials")
                        for i, ref in enumerate(st.session_state.references[:2]):  # Show only 2 most relevant references
                            st.markdown(f'<div class="reference-card">', unsafe_allow_html=True)
                            st.markdown(f'<div class="reference-title">{ref["title"]}</div>', unsafe_allow_html=True)
                            st.markdown(f'<div class="reference-snippet">{ref["snippet"]}</div>', unsafe_allow_html=True)
                            st.markdown(f'<a href="{ref["link"]}" target="_blank" class="reference-link">{ref["link"]}</a>', unsafe_allow_html=True)
                            st.markdown('</div>', unsafe_allow_html=True)
                
                with col2:
                    st.markdown(f"### {video['title']}")
                    st.markdown(f"**Channel:** {video['channel']}")
                    st.markdown(f"**Duration:** {video['duration']}")
                    st.markdown(f"**Views:** {video['views']:,}")
                    st.markdown(f"**URL:** [Open in YouTube]({video['url']})")
                    
                    # Related videos suggestion
                    st.markdown("### You might also like")
                    for i, related in enumerate(st.session_state.videos[:3]):
                        if related != video:
                            if st.button(f"▶️ {related['title'][:40]}...", key=f"related_{i}"):
                                st.session_state.selected_video = related
                                st.rerun()
                
                st.markdown('</div>', unsafe_allow_html=True)
  
        # Original Notes tab
        with tab7:
            if st.session_state.full_text:
                st.subheader("Original Lecture Notes")
                st.text_area("Full Notes Text", st.session_state.full_text, height=400)
            else:
                st.info("No notes text available.")
    else:
        # Welcome message when no document is uploaded
        st.markdown("""
        ## 🎓 Welcome to Lecture Notes Assistant Pro! 👋  
        Your all-in-one AI-powered study companion that transforms your lecture materials into powerful learning resources!  
        
        ### 🚀 What You Can Do:
        - **📖 Summarize** your lecture notes into concise, easy-to-read study guides.  
        - **🤖 AI Study Buddy** to answer your questions and help you understand complex topics.  
        - **❓ Generate FAQs** categorized by difficulty levels (Easy, Medium, Hard).  
        - **📝 Create Practice Questions** tailored to your course material.  
        - **🎴 Flashcards** with varying difficulty levels to reinforce learning.  
        - **🔗 Topic-Based Resources** – Get relevant website links & YouTube videos right inside the app!  
        
        ### 📌 Get Started in 3 Simple Steps:
        1️⃣ **Upload** your lecture notes (PDF, Word, PowerPoint, or Text).  
        2️⃣ **Enter** the subject or topic you want to study.  
        3️⃣ **Click** "Process Lecture Notes" and let AI handle the rest!  
        
        Need help preparing for an exam? Struggling with tough concepts? Your AI-powered study assistant is here to guide you every step of the way!  
        """)
        
        # Feature breakdown
        col1, col2 = st.columns(2)
        
        with col1:
            st.subheader("📚 Study Guide Features")
            st.markdown("""
            - **AI-Powered Summarization** – Uses BART and Gemini models to generate high-quality summaries.
            - **Personalized Content** – Tailored specifically for your subject and course material.
            - **Conversational & Easy-to-Understand** – No more robotic text!
            - **Feedback-Friendly** – Customize and refine your study guide.
            """)
            
            st.subheader("🎯 Question & Flashcard Generator")
            st.markdown("""
            - **Difficulty-Based Questions** – Automatically generates Easy, Medium, and Hard questions.
            - **Flashcards for Quick Learning** – Helps retain key concepts with a structured approach.
            - **Color-Coded Organization** – Makes learning more engaging.
            - **Exportable** – Download study materials in Markdown for offline review.
            """)
        
        with col2:
            st.subheader("🤖 AI Study Buddy")
            st.markdown("""
            - **Smart & Context-Aware** – Provides answers based on your specific lecture content.
            - **Interactive Q&A** – Ask follow-up questions and dive deeper into topics.
            - **Memory-Enhanced** – Remembers previous questions for seamless learning.
            """)
            
            st.subheader("🌎 External Learning Resources")
            st.markdown("""
            - **YouTube & Web Links** – Find curated videos and articles for deeper learning.
            - **Integrated Viewing** – Watch videos and access materials directly within the app.
            """)
            
            st.subheader("⚙️ Technical Highlights")
            st.markdown("""
            - **Supports Multiple File Formats** – Works with PDFs, DOCX, PPTX, and TXT.
            - **Optimized for Speed** – Handles large documents efficiently.
            - **NLP-Powered Concept Extraction** – Automatically identifies key points.
            - **State-of-the-Art AI Models** – Combines multiple AI techniques for superior results.
            """)

        st.markdown("""
        ---
        📌 **Supercharge Your Learning Experience Today!** Let AI help you study smarter, not harder. 🚀
        """)

# Run the app
if __name__ == "__main__":
    main()